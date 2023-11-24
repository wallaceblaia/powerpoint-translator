import sys
import os
from PyQt5.QtWidgets import QApplication, QSplashScreen, QMainWindow, QPushButton, QHBoxLayout, QVBoxLayout, QWidget, QFileDialog, QTextEdit, QComboBox, QProgressBar, QMessageBox, QAction, QDialog  # Adicionado QAction aqui
from PyQt5.QtCore import Qt, QTimer
from PyQt5.QtGui import QPixmap, QScreen
from pptx import Presentation
from googletrans import Translator
import requests
import re
import configparser
import six


class LicenseDialog(QDialog):
    def __init__(self, parent=None):
        super(LicenseDialog, self).__init__(parent)
        self.setWindowTitle("Termos de Licença")
        self.resize(500, 400)
        layout = QVBoxLayout(self)

        self.licenseText = QTextEdit()
        self.licenseText.setReadOnly(True)
        self.licenseText.setText(self.readLicenseText())
        layout.addWidget(self.licenseText)

        self.acceptButton = QPushButton("Aceitar", self)
        self.rejectButton = QPushButton("Recusar", self)
        self.acceptButton.clicked.connect(self.accept)
        self.rejectButton.clicked.connect(self.reject)

        layout.addWidget(self.acceptButton)
        layout.addWidget(self.rejectButton)

    def readLicenseText(self):
        try:
            with open(os.path.join('assets', 'licenca.txt'), 'r', encoding='utf-8') as file:
                return file.read()
        except FileNotFoundError:
            return "Não foi possível encontrar o arquivo de licença."




class TranslatorApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.licenseAccepted = False
        self.translator = Translator()
        self.from_code = "pt"
        self.to_code = "en"
        self.filename = ""
        self.output_filename = ""
                
    def initUI(self):
        # Configurações iniciais da janela
        self.setWindowTitle('ICM - PowerPoint Translator')
        self.setGeometry(100, 100, 800, 600)
        self.createMenuBar()

        # Layout principal
        mainLayout = QVBoxLayout()

        # Componentes da interface
        self.textEdit = QTextEdit()
        self.comboBox = QComboBox()
        self.languageDict = {"en": "English", "zh-cn": "Chinese", "hi": "Hindi", "es": "Spanish", "fr": "French"}
        for code, name in self.languageDict.items():
            self.comboBox.addItem(name, code)
        self.comboBox.currentIndexChanged.connect(self.languageChanged)
        self.progressBar = QProgressBar(self)
        self.selectFileBtn = QPushButton('Abrir Arquivo')
        self.translateBtn = QPushButton('Traduzir')
        self.saveBtn = QPushButton('Salvar Arquivo')
        self.saveBtn.setDisabled(True)

        # Layout dos botões
        buttonLayout = QHBoxLayout()
        buttonLayout.addWidget(self.selectFileBtn)
        buttonLayout.addWidget(self.translateBtn)
        buttonLayout.addWidget(self.saveBtn)

        # Adicionando componentes ao layout principal
        mainLayout.addWidget(self.textEdit)
        mainLayout.addWidget(self.comboBox)
        mainLayout.addWidget(self.progressBar)
        mainLayout.addLayout(buttonLayout)

        # Conectando sinais e slots
        self.selectFileBtn.clicked.connect(self.openFileDialog)
        self.translateBtn.clicked.connect(self.startTranslation)
        self.saveBtn.clicked.connect(self.saveFileDialog)

        # Configurando o widget central
        centralWidget = QWidget()
        centralWidget.setLayout(mainLayout)
        self.setCentralWidget(centralWidget)

        # Centralizando a janela
        qr = self.frameGeometry()
        cp = QScreen.availableGeometry(QApplication.primaryScreen()).center()
        qr.moveCenter(cp)
        self.move(qr.topLeft())

        # Exibindo a tela de licença
        self.showLicenseAgreement()
        self.show()

    def languageChanged(self, index):
        self.to_code = self.comboBox.itemData(index)
        
    def createMenuBar(self):
        menuBar = self.menuBar()
        aboutMenu = menuBar.addMenu('&Sobre')

        aboutAction = QAction('&Sobre', self)
        aboutAction.triggered.connect(self.showAboutDialog)
        aboutMenu.addAction(aboutAction)

    def showAboutDialog(self):
        QMessageBox.about(self, "Sobre ICM - PowerPoint Translator",
                          "ICM - PowerPoint Translator v1.0\n"
                          "Desenvolvido por Wallace Laia\n"
                          "E-mail: wallaceblaia@gmail.com\n\n"
                          "Este programa pode ser usado livremente para traduzir "
                          "arquivos do PowerPoint de slides criados pela Igreja Cristã Maranata. "
                          "O uso para outros fins viola os direitos de uso.")

    def showLicenseAgreement(self):
        config = configparser.ConfigParser()
        config_file = 'user_settings.ini'

        if not os.path.exists(config_file):
            dialog = LicenseDialog(self)
            if dialog.exec_() == QDialog.Accepted:
                config['Settings'] = {'LicenseAccepted': 'Yes'}
                with open(config_file, 'w') as configfile:
                    config.write(configfile)
                self.licenseAccepted = True
            else:
                sys.exit(0)
        else:
            config.read(config_file)
            self.licenseAccepted = config.getboolean('Settings', 'LicenseAccepted')

    def openFileDialog(self):
        self.filename, _ = QFileDialog.getOpenFileName(self, "Abra o Arquivo do PowerPoint ", "", "Arquivos do PowerPoint (*.pptx)")
        if self.filename:
            self.loadFileContent()

    def loadFileContent(self):
        prs = Presentation(self.filename)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        all_text.append(paragraph.text)
        self.textEdit.setText('\n'.join(all_text))

    def startTranslation(self):
        if self.filename:
            if self.check_internet():
                self.translate_file(self.filename)
            else:
                QMessageBox.warning(self, 'Error', 'Sem conexão com a Internet.')

    def saveFileDialog(self):
        if not self.output_filename or not os.path.exists(self.output_filename):
            QMessageBox.warning(self, 'Error', 'Não há arquivo traduzido para salvar.')
            return

        translated_filename, _ = QFileDialog.getSaveFileName(self, "Salve o arquivo traduzido", self.output_filename, "Arquivos do PowerPoint (*.pptx)")
        if translated_filename:
            os.rename(self.output_filename, translated_filename)
 
    def translate_file(self, filename):
        if not os.path.exists(filename):
            sys.exit("Arquivo de origem para tradução não encontrado")

        prs = Presentation(filename)
        total_slide_cnt = len(prs.slides)
        self.progressBar.setMaximum(total_slide_cnt)

        for i, slide in enumerate(prs.slides, start=1):
            if not self.check_internet():
                QMessageBox.warning(self, 'Error', 'Conexão com a Internet perdida. Por favor verifique sua conexão.')
                return

            for shape in slide.shapes:
                if shape.has_text_frame or shape.has_table or shape.shape_type == 6:
                    self.translate_shape(shape)
            self.progressBar.setValue(i)
            QApplication.processEvents()
        
        self.output_filename = filename.split('.')[0] + "_" + self.to_code + ".pptx"
        try:
            prs.save(self.output_filename)
            self.translationComplete()
        except Exception as e:
            QMessageBox.warning(self, 'Error', f'Failed to save the file: {str(e)}')

    def translationComplete(self):
        QMessageBox.information(self, 'Tradução Completa', 'A tradução foi concluída com sucesso!')
        self.saveBtn.setDisabled(False)
        self.openSaveDialog()

    def openSaveDialog(self):
        self.saveFileDialog()

    def check_internet(self):
        try:
            requests.get('http://www.google.com', timeout=3)
            return True
        except requests.ConnectionError:
            return False

    def translate_text(self, text):
        if isinstance(text, six.binary_type):
            text = text.decode("utf-8")

        sentences = re.split(r'(?<=[.!?]) +', text)
        translated_sentences = []

        for sentence in sentences:
            max_retries = 3
            for attempt in range(max_retries):
                try:
                    translated = self.translator.translate(sentence, src=self.from_code, dest=self.to_code).text
                    translated_sentences.append(translated)
                    break
                except Exception as e:
                    if attempt + 1 == max_retries:
                        translated_sentences.append(sentence)
                    continue

        return ' '.join(translated_sentences)

    def translate_paragraph(self, paragraph):
        if paragraph.text != '':
            result = self.translate_text(paragraph.text)

            p = paragraph._p
            for idx, run in enumerate(paragraph.runs):
                if idx == 0:
                    continue
                p.remove(run._r)

            if len(paragraph.runs) == 0:
                paragraph.add_run()

            font = paragraph.runs[0].font
            if font.name == 'Wingdings':
                font.name = 'Calibri'

            paragraph.runs[0].text = result

    def translate_shape(self, shape):
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                self.translate_paragraph(paragraph)
        elif shape.has_table:
            for cell in shape.table.iter_cells():
                if not cell.is_spanned and cell.text != '':
                    for paragraph in cell.text_frame.paragraphs:
                        self.translate_paragraph(paragraph)
        elif shape.shape_type == 6:
            for shp in shape.shapes:
                self.translate_shape(shp)
def postSplashInit(main_window, splash):
    splash.close()
    main_window.initUI() 
def main():
    app = QApplication(sys.argv)
    dir_path = os.path.dirname(os.path.realpath(__file__))
    splash_image_path = os.path.join(dir_path, "assets", "init.jpeg")
    splash_img = QPixmap(splash_image_path)
    scaled_splash_img = splash_img.scaled(600, 400, Qt.KeepAspectRatio, Qt.SmoothTransformation)
    splash = QSplashScreen(scaled_splash_img, Qt.WindowStaysOnTopHint)
    splash.show()

    main_window = TranslatorApp()
    QTimer.singleShot(3000, lambda: postSplashInit(main_window, splash))

    sys.exit(app.exec_())

if __name__ == '__main__':
    main()
